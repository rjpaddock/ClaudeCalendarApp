"""
Generate PowerPoint presentation documenting the Calendar Management application development
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "ASP.NET Core Calendar Management Application"
    subtitle.text = "Development Journey: From Concept to Implementation\nBuilt with Claude AI"

    # Slide 2: Initial Request
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Initial Request"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"Create an ASP.NET Core application using C# to manage calendars. This application should have a calendar view by month, week and day, should be able to manage users and create calendar invites for one or more users."'
    p.level = 1

    # Slide 3: Technical Decisions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Decisions"
    content = slide.placeholders[1].text_frame
    content.text = "Key Technology Choices:"

    decisions = [
        "Database: SQLite with Entity Framework Core",
        "UI Framework: ASP.NET Core with Razor Pages",
        "Authentication: None (focus on core functionality)",
        "Styling: Bootstrap 5",
        "Target Framework: .NET 9.0"
    ]
    for decision in decisions:
        p = content.add_paragraph()
        p.text = decision
        p.level = 1

    # Slide 4: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Database Architecture"
    content = slide.placeholders[1].text_frame
    content.text = "Three Core Tables:"

    tables = [
        "Users: Id, Name, Email, CreatedAt",
        "CalendarEvents: Id, Title, Description, StartDateTime, EndDateTime, Location, CreatedById",
        "EventAttendees: Id, EventId, UserId, ResponseStatus (Pending/Accepted/Declined/Tentative)"
    ]
    for table in tables:
        p = content.add_paragraph()
        p.text = table
        p.level = 1

    # Slide 5: Initial Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Initial Implementation"
    content = slide.placeholders[1].text_frame
    content.text = "Core Features Implemented:"

    features = [
        "Three calendar view types: Month, Week, Day",
        "User management (CRUD operations)",
        "Event creation with multi-user invitations",
        "Event details, editing, and deletion",
        "Attendee tracking with response status",
        "Navigation between different time periods"
    ]
    for feature in features:
        p = content.add_paragraph()
        p.text = feature
        p.level = 1

    # Slide 6: Iteration 1 - Remove Seed Data
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 1: Clean Start"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"Remove all seeded data from this application"'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Updated DbInitializer to start with empty database"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Removed all sample users and events"
    p.level = 1

    # Slide 7: Iteration 2 - Clickable Timeslots
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 2: Interactive Timeslots"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"When clicking on a timeslot open the create event window with the date and time populated."'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    changes = [
        "Added year/month/day/hour parameters to CreateEvent",
        "Made day numbers in Month view clickable",
        "Made timeslots in Week and Day views clickable",
        "Pre-populate event form with selected date/time"
    ]
    for change in changes:
        p = content.add_paragraph()
        p.text = change
        p.level = 1

    # Slide 8: Iteration 3 - Architecture Change
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 3: Major Architecture Change"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"Change all views to use Razor pages"'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Complete Refactor:"
    p.level = 0
    changes = [
        "Converted from MVC Controllers to Razor PageModels",
        "Updated Program.cs from AddControllersWithViews to AddRazorPages",
        "Created Pages folder structure with _Layout, _ViewStart, _ViewImports",
        "Updated all navigation from asp-action to asp-page",
        "Fixed build errors and updated routing"
    ]
    for change in changes:
        p = content.add_paragraph()
        p.text = change
        p.level = 1

    # Slide 9: Iteration 4 - Time Default
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 4: Better Time Defaults"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"When clicking on a day from the month or week calendar open the create event window with the date and time populated. If on a day, default the time to noon."'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Changed default time from 9 AM to 12 PM (noon)"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Maintained specific hour when clicking timeslots"
    p.level = 1

    # Slide 10: Iteration 5 - Enhanced Clickability
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 5: Enhanced User Experience"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"If a user clicks on any part of a day, not just the number, open the create event page"'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    changes = [
        "Added onclick handler to entire <td> day cell",
        "Added event.stopPropagation() to event items",
        "Added CSS hover effects for better UX",
        "Made entire day cell visually responsive"
    ]
    for change in changes:
        p = content.add_paragraph()
        p.text = change
        p.level = 1

    # Slide 11: Iteration 6 - Navigation Enhancement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 6: Visual Navigation"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"On week and month views add arrow icons to navigate by month and week respectively"'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Added ◄ and ► icons to Previous/Next buttons"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Applied to Month, Week, and Day views for consistency"
    p.level = 1

    # Slide 12: Iteration 7 - Layout Refinement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 7: Layout Refinement"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"Add the navigation buttons next to the description of the month and week"'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    changes = [
        "Moved navigation buttons inline with month/week title",
        "Used Bootstrap d-inline-block for side-by-side layout",
        "Maintained Create Event button on the right",
        "Improved visual hierarchy and space utilization"
    ]
    for change in changes:
        p = content.add_paragraph()
        p.text = change
        p.level = 1

    # Slide 13: Iteration 8 - Enhanced Navigation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 8: Enhanced Navigation"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"Move the navigation buttons back to the middle of the page and add 2 more buttons next and previous on the left and right of the month name or week name"'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    changes = [
        "Created git branch: feature/enhanced-navigation",
        "Added arrow-only buttons flanking the month/week title",
        "Restored full navigation button group to center",
        "Three-column layout: Title with arrows | Navigation | Create Event",
        "Improved accessibility with title attributes on arrow buttons"
    ]
    for change in changes:
        p = content.add_paragraph()
        p.text = change
        p.level = 1

    # Slide 14: Iteration 9 - Frontend/Backend Separation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 9: Frontend/Backend Separation"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"Rebuild this application. Separate all calls into a front end and a back end API."'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    changes = [
        "Created solution with two projects: CalendarManagement.API and CalendarManagement.Web",
        "API: ASP.NET Core with Controllers, EF Core, SQLite, and Swagger",
        "Web: ASP.NET Core Razor Pages calling API via HttpClient",
        "CalendarApiClient service for centralized API communication",
        "CORS configuration to allow cross-origin requests",
        "Separation of concerns: API handles data, Web handles presentation"
    ]
    for change in changes:
        p = content.add_paragraph()
        p.text = change
        p.level = 1

    # Slide 15: New Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "New Solution Architecture"
    content = slide.placeholders[1].text_frame
    content.text = "Two-Project Structure:"

    architecture = [
        "CalendarManagement.API (Backend):",
        "  - Models, Data, Controllers",
        "  - Entity Framework Core + SQLite",
        "  - RESTful API endpoints",
        "  - Swagger/OpenAPI documentation",
        "",
        "CalendarManagement.Web (Frontend):",
        "  - Razor Pages, ViewModels",
        "  - CalendarApiClient service",
        "  - HttpClient for API communication",
        "  - No direct database access"
    ]
    for item in architecture:
        if item:
            p = content.add_paragraph()
            p.text = item
            p.level = 0 if item.endswith(":") else 1
        else:
            content.add_paragraph()

    # Slide 16: File Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Structure"
    content = slide.placeholders[1].text_frame
    content.text = "Key Folders and Files:"

    structure = [
        "API: Models, Data (DbContext, DbInitializer), Controllers (Users, CalendarEvents)",
        "Web: Models (DTOs), ViewModels, Services (CalendarApiClient), Pages (Razor)",
        "Shared: wwwroot/css for styling",
        "Configuration: appsettings.json in both projects"
    ]
    for item in structure:
        p = content.add_paragraph()
        p.text = item
        p.level = 1

    # Slide 17: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Final Application Features"
    content = slide.placeholders[1].text_frame
    content.text = "What the Application Can Do:"

    features = [
        "View calendar in three formats: Month, Week, Day",
        "Navigate between time periods with arrow buttons",
        "Click anywhere on a day/timeslot to create events",
        "Create events with title, description, location, date/time",
        "Invite multiple users to events",
        "Track attendee response status (Pending/Accepted/Declined/Tentative)",
        "Manage users (Create, Read, Update, Delete)",
        "View event details with attendee list",
        "Edit and delete events"
    ]
    for feature in features:
        p = content.add_paragraph()
        p.text = feature
        p.level = 1

    # Slide 15: Technical Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Highlights"
    content = slide.placeholders[1].text_frame
    content.text = "Notable Implementation Details:"

    highlights = [
        "Entity Framework Core with SQLite for lightweight data storage",
        "Razor Pages architecture for cleaner separation of concerns",
        "Client-side JavaScript for interactive calendar cells",
        "event.stopPropagation() to prevent conflicting click handlers",
        "Bootstrap utility classes for responsive layout",
        "Custom CSS for calendar-specific styling and hover effects",
        "ViewModels for complex data binding scenarios",
        "Route parameters for deep linking to specific dates"
    ]
    for highlight in highlights:
        p = content.add_paragraph()
        p.text = highlight
        p.level = 1

    # Slide 19: Iteration 10 - Service Layer Refactoring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Iteration 10: Service Layer Refactoring"
    content = slide.placeholders[1].text_frame
    content.text = "User Request:"
    p = content.add_paragraph()
    p.text = '"Extract all logic from this application\'s controllers into a service layer"'
    p.level = 1

    p = content.add_paragraph()
    p.text = "Implementation:"
    p.level = 0
    changes = [
        "Created ICalendarEventService and CalendarEventService (10 methods, ~290 lines)",
        "Created IUserService and UserService (7 methods, ~95 lines)",
        "Registered services with dependency injection in Program.cs",
        "Refactored all 10 PageModels to use service layer",
        "Reduced PageModel code by ~70% on average",
        "Separated business logic from HTTP handling",
        "Improved testability, reusability, and maintainability"
    ]
    for change in changes:
        p = content.add_paragraph()
        p.text = change
        p.level = 1

    # Slide 20: Service Layer Benefits
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Service Layer Architecture Benefits"
    content = slide.placeholders[1].text_frame
    content.text = "Why Service Layer Matters:"

    benefits = [
        "Separation of Concerns: PageModels handle HTTP, services handle business logic",
        "Testability: Services can be unit tested independently with mocked data",
        "Reusability: Same services can power future API controllers",
        "Maintainability: Business logic changes in one place, not scattered across pages",
        "Scalability: Easy to add caching, logging, and other cross-cutting concerns",
        "Code Reduction: Calendar pages went from 60-70 lines to ~20 lines",
        "SOLID Principles: Follows Single Responsibility and Dependency Inversion"
    ]
    for benefit in benefits:
        p = content.add_paragraph()
        p.text = benefit
        p.level = 1

    # Slide 21: Lessons Learned
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Development Insights"
    content = slide.placeholders[1].text_frame
    content.text = "Key Takeaways:"

    lessons = [
        "Iterative development allows for rapid evolution",
        "User feedback drives feature refinement",
        "Architecture changes (MVC → Razor Pages → API/Frontend → Service Layer) are transformative",
        "Small UX improvements (clickable areas, icons) have big impact",
        "Starting with empty database is better for production readiness",
        "Interactive elements require careful event handling",
        "Visual cues (arrows, hover effects) improve usability",
        "Git branching enables safe experimentation with UI changes",
        "Separating frontend and backend enables scalability and maintainability",
        "Service layer refactoring dramatically improves code quality and maintainability"
    ]
    for lesson in lessons:
        p = content.add_paragraph()
        p.text = lesson
        p.level = 1

    # Slide 22: Technology Stack Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Complete Technology Stack"
    content = slide.placeholders[1].text_frame

    stack = [
        "Architecture: Service Layer Pattern with dependency injection",
        "Backend API: ASP.NET Core 9.0 with Controllers",
        "Frontend: ASP.NET Core 9.0 with Razor Pages",
        "Language: C# 12",
        "Database: SQLite",
        "ORM: Entity Framework Core 9.0",
        "API Documentation: Swagger/OpenAPI",
        "Communication: HttpClient, RESTful APIs",
        "UI Framework: Bootstrap 5",
        "Tools: Visual Studio / VS Code, dotnet CLI",
        "Version Control: Git"
    ]

    for tech in stack:
        p = content.add_paragraph()
        p.text = tech
        p.level = 0

    # Final Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    content = slide.placeholders[1].text_frame
    content.text = "Project Evolution:"

    p = content.add_paragraph()
    p.text = "From initial concept to fully functional application through 10 iterations"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Each user request refined and enhanced the application"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Demonstrates the power of iterative, AI-assisted development"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Final result: A professional, well-architected calendar management system with clean separation of concerns, ready for testing and deployment"
    p.level = 1

    # Save presentation
    output_path = "CalendarApp_Development_Journey.pptx"
    prs.save(output_path)
    print(f"Presentation saved as: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
